# -*- coding: utf-8 -*-
"""사진 버전 빌드 : 1~4번 박스에 사진 3장씩(총 12장)을 타일로 삽입.

photos/box{n}_{1..3}.(jpg|jpeg|png|webp) 가 있으면 그 사진을,
없으면 assets/photo_placeholders/ 의 자리표시자를 넣는다.
원본 비율과 무관하게 타일 비율(약 3:2)로 가운데 크롭되며,
각 사진은 PowerPoint에서 [그림 바꾸기]로 개별 교체할 수 있다.
"""
import os, sys
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)
from sdl_layout import (apply_geometry, emu, BAND_X, BAND_TOP, IMG_DY,
                        TILE_PAD, TILE_GAP, TILE_W, TILE_H, SUBJECTS, ALT, NS_A)

SRC   = os.path.join(ROOT, 'assets', '원본_v4_new_mod.pptx')
OUT   = os.path.join(ROOT, '자율실험실_SDL_도식_v5-photo_사진버전.pptx')
PHOTO = os.path.join(ROOT, 'photos')
PLACE = os.path.join(HERE, 'photo_placeholders')
EXTS  = ('.jpg', '.jpeg', '.png', '.webp')
TARGET_AR = TILE_W / TILE_H


def find_photo(n, i):
    for e in EXTS:
        p = os.path.join(PHOTO, f'box{n}_{i}{e}')
        if os.path.exists(p):
            return p, True
    return os.path.join(PLACE, f'box{n}_{i}.png'), False


def center_crop(pic, path):
    """원본 비율을 타일 비율에 맞춰 가운데 크롭 (srcRect)."""
    with Image.open(path) as im:
        ar = im.width / im.height
    if abs(ar - TARGET_AR) < 1e-4:
        return
    if ar > TARGET_AR:                       # 원본이 더 넓다 -> 좌우 잘라냄
        f = (1 - TARGET_AR / ar) / 2
        pic.crop_left = pic.crop_right = f
    else:                                    # 원본이 더 높다 -> 위아래 잘라냄
        f = (1 - ar / TARGET_AR) / 2
        pic.crop_top = pic.crop_bottom = f


def style(pic):
    """모서리 둥글게 + 얇은 테두리 (카드 디자인과 통일)."""
    spPr = pic._element.spPr
    old = spPr.find(qn('a:prstGeom'))
    if old is not None:
        spPr.remove(old)
    geom = parse_xml(
        f'<a:prstGeom xmlns:a="{NS_A}" prst="roundRect">'
        f'<a:avLst><a:gd name="adj" fmla="val 7000"/></a:avLst></a:prstGeom>')
    xfrm = spPr.find(qn('a:xfrm'))
    xfrm.addnext(geom)
    geom.addnext(parse_xml(
        f'<a:ln xmlns:a="{NS_A}" w="12700"><a:solidFill>'
        f'<a:srgbClr val="DDE0E6"/></a:solidFill></a:ln>'))


def main():
    prs = Presentation(SRC)
    slide = prs.slides[0]
    apply_geometry(slide)

    missing = []
    for n in (1, 2, 3, 4):
        y = emu(BAND_TOP[n] + IMG_DY) + emu(TILE_PAD)
        for i in range(3):
            path, real = find_photo(n, i + 1)
            if not real:
                missing.append(f'photos/box{n}_{i+1}.jpg  ({SUBJECTS[n][i]})')
            x = BAND_X[n] + emu(TILE_PAD) + i * (TILE_W + emu(TILE_GAP))
            pic = slide.shapes.add_picture(path, Emu(x), Emu(y), Emu(TILE_W), Emu(TILE_H))
            pic.name = f'예시 사진 {n}-{i+1} {SUBJECTS[n][i]}'
            pic._element.find('.//' + qn('p:cNvPr')).set(
                'descr', f'{ALT[n]} — {SUBJECTS[n][i]}')
            center_crop(pic, path)
            style(pic)

    prs.save(OUT)
    print('saved', os.path.basename(OUT))
    if missing:
        print(f'\n자리표시자로 채운 슬롯 {len(missing)}개 (실제 사진을 넣으면 자동 대체):')
        for m in missing:
            print('  -', m)


if __name__ == '__main__':
    main()
