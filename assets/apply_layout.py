# -*- coding: utf-8 -*-
"""1~4번 박스에 예시 이미지 삽입 + 이미지 공간 확보를 위한 도식 레이아웃 재조정."""
from pptx import Presentation
from pptx.util import Emu
import copy

E = 914400
def emu(v): return int(round(v * E))

prs = Presentation('assets/원본_v4_new_mod.pptx')
slide = prs.slides[0]

# ---- id -> element 인덱스 (그룹 내부 포함하지 않음: 최상위만) ----
tree = slide.shapes._spTree
ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
      'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
elems = {}
for el in tree:
    tag = el.tag.split('}')[1]
    if tag not in ('sp', 'pic', 'grpSp', 'cxnSp'):
        continue
    cNvPr = el.find('.//{%s}cNvPr' % ns['p'])
    elems[int(cNvPr.get('id'))] = el

def xfrm_of(el):
    tag = el.tag.split('}')[1]
    holder = el.find('{%s}spPr' % ns['p'])
    if holder is None:
        holder = el.find('{%s}grpSpPr' % ns['p'])
    return holder.find('{%s}xfrm' % ns['a'])

def place(sid, y=None, h=None):
    xf = xfrm_of(elems[sid])
    if y is not None:
        xf.find('{%s}off' % ns['a']).set('y', str(emu(y)))
    if h is not None:
        xf.find('{%s}ext' % ns['a']).set('cy', str(emu(h)))

# =========================== 세로 레이아웃 재배치 ===========================
R1, R2, CARD_H = 3.122, 5.934, 2.372          # 카드 행 위치 / 높이(1.62" -> 2.372")
IMG_DY, IMG_H  = 1.64, 0.62                   # 카드 상단 기준 이미지 밴드 위치/높이

layout = {
    # 상단부 (여백·박스 높이를 줄여 아래쪽 공간 확보)
    2: (0.22, 0.64), 3: (0.22, 0.64), 4: (0.328, None),
    5: (0.98, 0.64), 6: (0.98, 0.64), 7: (1.046, None),
    8: (1.74, 0.87), 92: (1.7527, None), 59: (1.9015, None),
    9: (2.65, None),
    # 자율실험 폐루프 컨테이너 + 라벨
    10: (2.922, 6.4568), 11: (2.752, None), 85: (2.7537, None),
    # 1~4번 카드
    12: (R1, CARD_H), 13: (R1 + 0.13, None), 15: (R1 + 0.1521, None),
    16: (R1 + 0.4999, None), 17: (R1 + 0.58, None), 18: (R1 + 1.15, None),
    19: (R1, CARD_H), 20: (R1 + 0.13, None), 22: (R1 + 0.1521, None),
    23: (R1 + 0.4999, None), 24: (R1 + 0.58, None), 25: (R1 + 1.15, None),
    26: (R2, CARD_H), 27: (R2 + 0.13, None), 29: (R2 + 0.0343, None),
    30: (R2 + 0.4999, None), 31: (R2 + 0.58, None), 32: (R2 + 1.0670, None),
    33: (R2, CARD_H), 34: (R2 + 0.13, None), 36: (R2 + 0.1521, None),
    37: (R2 + 0.4999, None), 38: (R2 + 0.58, None), 39: (R2 + 1.0670, None),
    # 순환 화살표
    40: (4.198, None), 42: (7.0095, None), 41: (5.554, None), 43: (5.554, None),
    # 하단 지원 레이어
    54: (8.426, None), 55: (8.436, None),
    # 우측 지원 요소(그룹) + 헤더 텍스트
    71: (3.533, None), 87: (3.5464, None),
    72: (6.345, None), 89: (6.3631, None),
    # 연결선
    69: (4.2403, None), 80: (4.3514, 2.8104), 74: (7.2568, None), 66: (7.120, 2.8638),
    # 최하단 비교 박스
    61: (9.5588, 0.85), 56: (9.6203, None),
}
for sid, (y, h) in layout.items():
    place(sid, y, h)

# 1번 카드의 ex. 텍스트 상자만 폭이 좁아(2.33") 줄바꿈이 과하게 발생 -> 나머지 카드와 동일한 폭으로 정렬
xf = xfrm_of(elems[18])
xf.find('{%s}off' % ns['a']).set('x', '809468')
xf.find('{%s}ext' % ns['a']).set('cx', '2400337')

# =========================== 예시 이미지 삽입 ===========================
IMG_W = 2400337                                  # 카드 내부 구분선과 동일 폭
CARDS = [(1, 809468,  R1, 'AI 실험 설계 예시: AI 추천엔진(신경망) · 과거 실험 데이터베이스 · 가상실험 모델'),
         (2, 4368927, R1, '로봇·장비 실험 수행 예시: 로봇팔 · 액체 핸들러와 웰플레이트 · 측정장비'),
         (3, 4368927, R2, 'AI·분석SW 결과분석 예시: 측정데이터 곡선 분석 · 목표 충족여부 자동판정 · 결과 시각화 대시보드'),
         (4, 809468,  R2, 'AI 모델 갱신 예시: 모델 재학습 · 파인튜닝 파이프라인 · 다음 실험조건 또는 종료 결정')]

for n, x, top, alt in CARDS:
    pic = slide.shapes.add_picture(f'assets/example_images/box{n}.png', Emu(x), Emu(emu(top + IMG_DY)),
                                   Emu(IMG_W), Emu(emu(IMG_H)))
    pic.name = f'예시 이미지 {n}'
    pic._element.find('.//{%s}cNvPr' % ns['p']).set('descr', alt)

prs.save('자율실험실_SDL_도식_v5_예시이미지.pptx')
print('saved 자율실험실_SDL_도식_v5_예시이미지.pptx')
